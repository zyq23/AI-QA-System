from __future__ import annotations

from pathlib import Path

from docx import Document
from PIL import Image, ImageDraw
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont

from app.parsers.service import DocumentParserService
from app.parsers.pdf_parser import PaddleOcrAdapter
from app.parsers.ocr_utils import clean_ocr_text, split_visual_text
from app.parsers.pptx_parser import _collapse_ocr_heavy_slide, _compact_ocr_parts, _is_placeholder_text, _merge_dense_ocr_elements
from app.services.chunker import ChunkerService, ChunkingContext
from scripts.inspect_ppt_parse import build_slide_summaries


def build_docx(path: Path) -> None:
    document = Document()
    document.add_heading("G1 产品手册", level=1)
    document.add_paragraph("G1 支持本地资料问答。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "参数"
    table.cell(0, 1).text = "值"
    table.cell(1, 0).text = "关节"
    table.cell(1, 1).text = "多自由度"
    document.save(path)


def build_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "G1 方案概览"
    box = slide.shapes.add_textbox(left=1_000_000, top=2_000_000, width=6_000_000, height=2_000_000)
    box.text_frame.text = "支持引用式回答与知识库检索"
    prs.save(path)


def build_png(path: Path) -> None:
    image = Image.new("RGB", (240, 80), color="white")
    drawer = ImageDraw.Draw(image)
    drawer.text((12, 24), "OCR图片知识点", fill="black")
    image.save(path)


def build_pdf(path: Path) -> None:
    pdf = canvas.Canvas(str(path))
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    pdf.setFont("STSong-Light", 14)
    pdf.drawString(72, 760, "G1 FAQ")
    pdf.drawString(72, 730, "资料不足时需要明确提示知识库没有直接证据。")
    pdf.save()


def test_parse_docx_pptx_pdf(tmp_path: Path):
    parser = DocumentParserService(enable_ocr_fallback=False)
    docx_path = tmp_path / "sample.docx"
    pptx_path = tmp_path / "sample.pptx"
    pdf_path = tmp_path / "sample.pdf"
    build_docx(docx_path)
    build_pptx(pptx_path)
    build_pdf(pdf_path)

    docx_doc = parser.parse(docx_path)
    pptx_doc = parser.parse(pptx_path)
    pdf_doc = parser.parse(pdf_path)

    assert docx_doc.blocks
    assert "G1 产品手册" in docx_doc.raw_markdown
    assert any("多自由度" in block.content for block in docx_doc.blocks)

    assert pptx_doc.blocks
    assert "G1 方案概览" in pptx_doc.raw_markdown

    assert pdf_doc.blocks
    assert any("知识库没有直接证据" in block.content for block in pdf_doc.blocks)


def test_parse_pptx_image_ocr(tmp_path: Path, monkeypatch):
    parser = DocumentParserService(enable_ocr_fallback=True)
    pptx_path = tmp_path / "image_ocr.pptx"
    png_path = tmp_path / "ocr.png"
    build_png(png_path)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "图片知识点"
    slide.shapes.add_picture(str(png_path), left=1_000_000, top=1_500_000, width=2_400_000, height=800_000)
    prs.save(pptx_path)

    monkeypatch.setattr(
        parser.pptx_parser.ocr,
        "extract_image_text",
        lambda raw: "这段内容来自图片OCR",
    )

    pptx_doc = parser.parse(pptx_path)

    assert pptx_doc.ocr_used is True
    assert any(block.kind == "image_ocr" and "图片OCR" in block.content for block in pptx_doc.blocks)
    assert "这段内容来自图片OCR" in pptx_doc.raw_markdown


def test_clean_ocr_text_filters_noise_but_keeps_key_facts():
    cleaned, quality = clean_ocr_text(
        "金1 银合金！！ 错会金1 ???\n"
        "Unitree G1 总自由度约43个\n"
        "适用课程：Python程序设计、深度学习\n"
    )

    assert "Unitree G1 总自由度约43个" in cleaned
    assert "适用课程：Python程序设计、深度学习" in cleaned
    assert "错会金1" not in cleaned
    assert quality > 0.45


def test_clean_ocr_text_salvages_low_confidence_but_meaningful_lines():
    cleaned, quality = clean_ocr_text(
        "产业学院治理模式\n"
        "协同育人机制\n"
        "???!\n"
    )

    assert "产业学院治理模式" in cleaned
    assert "协同育人机制" in cleaned
    assert quality >= 0.18


def test_split_visual_text_merges_short_fragmented_chinese_parts():
    parts = split_visual_text("已\n出\n版\n2024年轩辕网络将联合全国多所高校专业教师")

    assert parts[0] == "已出版"
    assert "2024年轩辕网络将联合全国多所高校专业教师" in parts


def test_parse_pptx_image_ocr_marks_low_confidence_blocks(tmp_path: Path, monkeypatch):
    parser = DocumentParserService(enable_ocr_fallback=True)
    pptx_path = tmp_path / "image_ocr_low_conf.pptx"
    png_path = tmp_path / "ocr_low_conf.png"
    build_png(png_path)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "低置信图片知识点"
    slide.shapes.add_picture(str(png_path), left=1_000_000, top=1_500_000, width=2_400_000, height=800_000)
    prs.save(pptx_path)

    monkeypatch.setattr(
        parser.pptx_parser.ocr,
        "extract_image_text",
        lambda raw: "金1 银合金！！ 错会金1 ???\nUnitree G1 总自由度约43个",
    )

    pptx_doc = parser.parse(pptx_path)

    low_conf_blocks = [block for block in pptx_doc.blocks if block.kind == "image_ocr_low_conf"]
    assert low_conf_blocks
    assert any("总自由度约43个" in block.content for block in low_conf_blocks)
    assert all(block.quality_score < 0.62 for block in low_conf_blocks)


def test_parse_pptx_tolerates_shapes_with_none_coordinates(tmp_path: Path, monkeypatch):
    parser = DocumentParserService(enable_ocr_fallback=False)
    pptx_path = tmp_path / "none_coords.pptx"
    build_pptx(pptx_path)

    class FakeParagraph:
        def __init__(self, text: str):
            self.runs = []
            self.text = text

    class FakeTextFrame:
        def __init__(self, text: str):
            self.paragraphs = [FakeParagraph(text)]

    class FakeShape:
        def __init__(self):
            self.top = None
            self.left = None
            self.shape_type = None
            self.has_text_frame = True
            self.text = "占位内容"
            self.text_frame = FakeTextFrame("占位内容")
            self.has_table = False

    class FakeSlideShapes(list):
        title = None

    class FakeSlide:
        def __init__(self):
            self.shapes = FakeSlideShapes([FakeShape()])

    class FakePresentation:
        def __init__(self):
            self.slides = [FakeSlide()]

    monkeypatch.setattr("app.parsers.pptx_parser.Presentation", lambda _: FakePresentation())

    parsed = parser.parse(pptx_path)

    assert parsed.blocks
    assert parsed.blocks[0].content == "占位内容"
    assert "占位内容" in parsed.raw_markdown


def test_parse_pptx_filters_placeholder_template_text(tmp_path: Path, monkeypatch):
    parser = DocumentParserService(enable_ocr_fallback=False)
    pptx_path = tmp_path / "placeholder_text.pptx"
    build_pptx(pptx_path)

    class FakeParagraph:
        def __init__(self, text: str):
            self.runs = []
            self.text = text

    class FakeTextFrame:
        def __init__(self, *texts: str):
            self.paragraphs = [FakeParagraph(text) for text in texts]

    class FakeShape:
        def __init__(self):
            self.top = 0
            self.left = 0
            self.shape_type = None
            self.has_text_frame = True
            self.text = "Click to add text"
            self.text_frame = FakeTextFrame(
                "Click to add text",
                "Lorem commodo et tation consetetur esse eirmod sed lorem.",
                "真实内容保留",
            )
            self.has_table = False

    class FakeSlideShapes(list):
        title = None

    class FakeSlide:
        def __init__(self):
            self.shapes = FakeSlideShapes([FakeShape()])

    class FakePresentation:
        def __init__(self):
            self.slides = [FakeSlide()]

    monkeypatch.setattr("app.parsers.pptx_parser.Presentation", lambda _: FakePresentation())

    parsed = parser.parse(pptx_path)

    contents = [block.content for block in parsed.blocks]
    assert "真实内容保留" in contents
    assert all("Click to add text" not in content for content in contents)
    assert all("Lorem commodo" not in content for content in contents)


def test_placeholder_filter_catches_more_template_variants():
    assert _is_placeholder_text("Double click to edit")
    assert _is_placeholder_text("Replace me with your content")
    assert _is_placeholder_text("Click to add table")
    assert not _is_placeholder_text("联合培养机制与课程体系")


def test_compact_ocr_parts_merges_dense_short_fragments():
    parts = [
        "大数系列丛书",
        "大数据技术基础",
        "应用教程",
        "周张线填/苏、邱新羽主演",
        "清华大学出版社",
        "程序设计基础",
        "中国工信出版社",
        "人民邮电出版社",
        "申时全上",
        "大数据开发技术",
    ]

    compacted = _compact_ocr_parts(parts, 0.96)

    assert len(compacted) < len(parts)
    assert any("大数系列丛书" in part for part in compacted)
    assert any("大数据开发技术" in part for part in compacted)


def test_compact_ocr_parts_drops_isolated_numeric_noise_when_substantive_text_exists():
    parts = ["10", "1112", "234", "研究生联合培养基地", "校企合作AI创新实验班", "8", "7"]

    compacted = _compact_ocr_parts(parts, 0.74)

    assert all(part not in compacted for part in ["10", "1112", "234", "8", "7"])
    assert any("研究生联合培养基地" in part for part in compacted)
    assert any("校企合作AI创新实验班" in part for part in compacted)


def test_merge_dense_ocr_elements_coalesces_picture_wall_fragments():
    elements = [
        (1000000, 1000000, "大数系列丛书", "image_ocr", "picture-ocr-1", 0.96),
        (1020000, 1500000, "大数据技术基础", "image_ocr", "picture-ocr-2", 0.96),
        (1050000, 2000000, "应用教程", "image_ocr", "picture-ocr-3", 0.96),
        (1100000, 2600000, "程序设计基础", "image_ocr", "picture-ocr-4", 0.96),
        (1150000, 3200000, "人民邮电出版社", "image_ocr", "picture-ocr-5", 0.96),
        (1400000, 1000000, "清华大学出版社", "image_ocr", "picture-ocr-6", 0.96),
        (1420000, 1500000, "大数据开发技术", "image_ocr", "picture-ocr-7", 0.96),
        (1450000, 2000000, "数据分析基础", "image_ocr", "picture-ocr-8", 0.96),
        (1480000, 2600000, "算法与实现", "image_ocr", "picture-ocr-9", 0.96),
        (1500000, 3200000, "教材案例", "image_ocr", "picture-ocr-10", 0.96),
        (1750000, 1000000, "更多教材", "image_ocr", "picture-ocr-11", 0.96),
        (1780000, 1500000, "专题实践", "image_ocr", "picture-ocr-12", 0.96),
        (1800000, 2000000, "创新实验", "image_ocr", "picture-ocr-13", 0.96),
        (1830000, 2600000, "在线精品", "image_ocr", "picture-ocr-14", 0.96),
        (1850000, 3200000, "课程资源", "image_ocr", "picture-ocr-15", 0.96),
        (2100000, 1000000, "配套资料", "image_ocr", "picture-ocr-16", 0.96),
        (2120000, 1500000, "案例库", "image_ocr", "picture-ocr-17", 0.96),
        (2150000, 2000000, "实践题库", "image_ocr", "picture-ocr-18", 0.96),
        (2170000, 2600000, "实训指引", "image_ocr", "picture-ocr-19", 0.96),
        (2200000, 3200000, "教学成果", "image_ocr", "picture-ocr-20", 0.96),
    ]

    merged = _merge_dense_ocr_elements(elements)

    assert len(merged) < len(elements)
    assert any(item[4] == "picture-ocr-merged" for item in merged)
    assert any("大数系列丛书" in item[2] for item in merged)
    assert any("教学成果" in item[2] for item in merged)


def test_merge_dense_ocr_elements_keeps_normal_slide_when_ocr_is_sparse():
    elements = [
        (1000000, 1000000, "AI赋能阶段", "paragraph", "body", 1.0),
        (1200000, 1000000, "2024-2025年，提升教育服务的智能化水平", "paragraph", "body", 1.0),
        (1500000, 1500000, "广东技术师范大学-广东轩辕网络科技股份有限公司", "image_ocr", "picture-ocr-1", 0.92),
        (1800000, 1500000, "研究生联合培养基地揭牌仪式", "image_ocr", "picture-ocr-2", 0.92),
    ]

    merged = _merge_dense_ocr_elements(elements)

    assert merged == elements


def test_collapse_ocr_heavy_slide_compacts_picture_wall():
    ocr_elements = [
        (1_000_000 + i * 80_000, 1_000_000 + (i % 5) * 700_000, f"教材条目{i}", "image_ocr", f"picture-ocr-{i}", 0.95)
        for i in range(60)
    ]
    elements = [(900_000, 800_000, "教材建设服务", "paragraph", "body", 1.0)] + ocr_elements

    collapsed = _collapse_ocr_heavy_slide(elements)

    assert len(collapsed) < len(elements)
    assert any(item[4] == "picture-ocr-collapsed" for item in collapsed)
    assert any("教材条目0" in item[2] for item in collapsed if item[4] == "picture-ocr-collapsed")


def test_parse_pptx_prefers_column_order_when_slide_has_two_clear_columns(tmp_path: Path, monkeypatch):
    parser = DocumentParserService(enable_ocr_fallback=False)
    pptx_path = tmp_path / "two_columns.pptx"
    build_pptx(pptx_path)

    class FakeParagraph:
        def __init__(self, text: str):
            self.runs = []
            self.text = text

    class FakeTextFrame:
        def __init__(self, text: str):
            self.paragraphs = [FakeParagraph(text)]

    class FakeShape:
        def __init__(self, text: str, top: int, left: int):
            self.top = top
            self.left = left
            self.shape_type = None
            self.has_text_frame = True
            self.text = text
            self.text_frame = FakeTextFrame(text)
            self.has_table = False

    class FakeSlideShapes(list):
        title = None

    class FakeSlide:
        def __init__(self):
            self.shapes = FakeSlideShapes(
                [
                    FakeShape("左列上", 1_000_000, 1_000_000),
                    FakeShape("右列上", 1_100_000, 5_000_000),
                    FakeShape("左列下", 2_000_000, 1_000_000),
                    FakeShape("右列下", 2_100_000, 5_000_000),
                ]
            )

    class FakePresentation:
        def __init__(self):
            self.slides = [FakeSlide()]

    monkeypatch.setattr("app.parsers.pptx_parser.Presentation", lambda _: FakePresentation())

    parsed = parser.parse(pptx_path)

    contents = [block.content for block in parsed.blocks]
    assert contents == ["左列上", "左列下", "右列上", "右列下"]


def test_parse_pdf_splits_blocks_and_marks_table_like_content(tmp_path: Path):
    parser = DocumentParserService(enable_ocr_fallback=False)
    pdf_path = tmp_path / "structured.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    pdf.setFont("STSong-Light", 14)
    pdf.drawString(72, 760, "申请步骤")
    pdf.drawString(72, 738, "1. 了解项目要求")
    pdf.drawString(72, 716, "2. 提交相关申请")
    pdf.drawString(72, 694, "参数    数值    说明")
    pdf.drawString(72, 672, "平台    官网    在线申请")
    pdf.save()

    parsed = parser.parse(pdf_path)

    assert any("了解项目要求" in block.content for block in parsed.blocks)
    assert any(block.kind == "table" and "参数 数值 说明" in block.content for block in parsed.blocks)
    assert any("文本" in block.section_path or "区块" in block.section_path for block in parsed.blocks)


def test_chunker_keeps_low_confidence_ocr_more_compact():
    parsed = parser = None
    parsed = type(
        "Parsed",
        (),
        {
            "blocks": [
                type("Block", (), {"page_or_slide": "slide-1", "section_path": "A", "content": "第一句。第二句。第三句。第四句。", "kind": "paragraph", "quality_score": 1.0})(),
                type("Block", (), {"page_or_slide": "slide-1", "section_path": "A", "content": "OCR一行\nOCR二行\nOCR三行", "kind": "image_ocr_low_conf", "quality_score": 0.4})(),
            ]
        },
    )()
    service = ChunkerService(target_size=20, overlap=0)
    chunks = service.chunk(
        parsed,
        ChunkingContext(
            document_id="doc-1",
            version_id="ver-1",
            file_name="sample.pptx",
            source_type="upload",
            trust_level="internal",
            target_size=20,
            overlap=0,
        ),
    )

    assert len(chunks) >= 3
    assert any("OCR一行" in chunk.plain_text for chunk in chunks)


def test_chunker_drops_short_low_confidence_ocr_noise():
    parsed = type(
        "Parsed",
        (),
        {
            "blocks": [
                type("Block", (), {"page_or_slide": "slide-1", "section_path": "A", "content": "1112\n234\n8\n研究生联合培养基地", "kind": "image_ocr_low_conf", "quality_score": 0.3})(),
            ]
        },
    )()
    service = ChunkerService(target_size=20, overlap=0)
    chunks = service.chunk(
        parsed,
        ChunkingContext(
            document_id="doc-1",
            version_id="ver-1",
            file_name="sample.pptx",
            source_type="upload",
            trust_level="internal",
            target_size=20,
            overlap=0,
        ),
    )

    assert chunks
    assert any("研究生联合培养基地" in chunk.plain_text for chunk in chunks)
    assert all("1112" not in chunk.plain_text for chunk in chunks)
    assert all("234" not in chunk.plain_text for chunk in chunks)


def test_slide_summary_counts_placeholder_noise_and_chunks():
    parsed = type(
        "Parsed",
        (),
        {
            "blocks": [
                type("Block", (), {"page_or_slide": "slide-79", "section_path": "A", "content": "Click to add text", "kind": "paragraph", "quality_score": 1.0})(),
                type("Block", (), {"page_or_slide": "slide-79", "section_path": "A", "content": "123", "kind": "image_ocr_low_conf", "quality_score": 0.2})(),
                type("Block", (), {"page_or_slide": "slide-79", "section_path": "A", "content": "真实内容保留", "kind": "paragraph", "quality_score": 1.0})(),
            ],
            "warnings": ["第 79 页图片 OCR 失败: mock"],
        },
    )()

    summaries = build_slide_summaries(parsed, {"slide-79": 2})

    assert len(summaries) == 1
    summary = summaries[0]
    assert summary.slide == "slide-79"
    assert summary.warning_count == 1
    assert summary.placeholder_hits == 1
    assert summary.low_quality_short_noise_hits == 1
    assert summary.chunk_count == 2


def test_extract_image_text_prefers_rapidocr_when_paddle_is_unstable(monkeypatch):
    adapter = PaddleOcrAdapter(language="ch")

    class FakeRapid:
        def __call__(self, image):
            return [[None, "产业学院治理模式", 0.95], [None, "协同育人机制", 0.88]], None

    monkeypatch.setattr(adapter, "_load_rapid", lambda: FakeRapid())
    monkeypatch.setattr(adapter, "extract_array_text", lambda image: (_ for _ in ()).throw(RuntimeError("should not fallback")))

    image = Image.new("RGB", (200, 60), color="white")
    from io import BytesIO

    payload = BytesIO()
    image.save(payload, format="PNG")

    text = adapter.extract_image_text(payload.getvalue())

    assert "产业学院治理模式" in text
    assert "协同育人机制" in text


def test_parse_pdf_extracts_embedded_image_ocr(tmp_path: Path, monkeypatch):
    parser = DocumentParserService(enable_ocr_fallback=True)
    pdf_path = tmp_path / "image_ocr.pdf"
    png_path = tmp_path / "pdf_ocr.png"
    build_png(png_path)

    pdf = canvas.Canvas(str(pdf_path))
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    pdf.setFont("STSong-Light", 14)
    pdf.drawString(72, 760, "图文页")
    pdf.drawImage(str(png_path), 72, 640, width=180, height=60)
    pdf.save()

    monkeypatch.setattr(
        parser.pdf_parser.ocr,
        "extract_image_text",
        lambda raw: "华为通过理论重构、架构重构、软件重构实现突围",
    )

    parsed = parser.parse(pdf_path)

    assert parsed.ocr_used is True
    assert any(block.kind in {"image_ocr", "image_ocr_low_conf"} and "理论重构" in block.content for block in parsed.blocks)
