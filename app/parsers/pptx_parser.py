from __future__ import annotations

import re
from pathlib import Path
from typing import Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.domain import ParsedDocument, SourceBlock
from app.parsers.ocr_utils import LIST_MARKER_RE, clean_ocr_text, split_visual_text
from app.parsers.pdf_parser import PaddleOcrAdapter
from app.utils import normalize_text

PPT_PLACEHOLDER_PREFIXES = (
    "click to add ",
    "lorem ",
    "lorem ipsum",
    "double click to edit",
    "replace me",
)

PPT_PLACEHOLDER_PATTERNS = (
    re.compile(r"^click to add (?:title|subtitle|text|notes|chart|table|picture|media|content)\b", re.IGNORECASE),
    re.compile(r"^double click to edit\b", re.IGNORECASE),
    re.compile(r"^replace me\b", re.IGNORECASE),
    re.compile(r"^lorem(?: ipsum)?\b", re.IGNORECASE),
)


def _iter_flattened_shapes(shapes) -> Iterator:
    """Recursively yield all leaf shapes from a shape tree, unwrapping GROUP containers."""
    for shape in shapes:
        shape_type = getattr(shape, "shape_type", None)
        if shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            yield from _iter_flattened_shapes(shape.shapes)
        else:
            yield shape


def _safe_coord(shape, attr: str) -> int:
    value = getattr(shape, attr, 0)
    if value is None:
        return 0
    try:
        return int(value)
    except (TypeError, ValueError):
        return 0


def _table_to_markdown(table) -> str:
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [normalize_text(cell.text) for cell in row.cells]
        if any(cells):
            rows.append(cells)
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    padded = [row + [""] * (width - len(row)) for row in rows]
    header = padded[0]
    divider = ["---"] * width
    body = padded[1:] or [[""] * width]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(divider) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _shape_role(shape, slide_title: str) -> str:
    text = ""
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        text = normalize_text(shape.text)
    if slide_title and text and text == slide_title:
        return "title"
    if hasattr(shape, "has_table") and shape.has_table:
        return "table"
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if LIST_MARKER_RE.match(text):
        return "list"
    return "body"


def _split_shape_text(text: str) -> list[str]:
    return split_visual_text(text)


def _is_substantive_ocr_part(text: str) -> bool:
    normalized = normalize_text(text)
    if not normalized:
        return False
    if len(normalized) >= 5 and re.search(r"[\u4e00-\u9fffA-Za-z]{2,}", normalized):
        return True
    if re.search(r"[\u4e00-\u9fffA-Za-z]{3,}", normalized):
        return True
    if re.search(r"\d{4,}", normalized) and re.search(r"[\u4e00-\u9fffA-Za-z]", normalized):
        return True
    return False


def _compact_ocr_parts(parts: list[str], quality_score: float) -> list[str]:
    normalized_parts = [normalize_text(part) for part in parts if normalize_text(part)]
    if not normalized_parts:
        return []

    deduped_parts: list[str] = []
    seen_recent: set[str] = set()
    for part in normalized_parts:
        if part in seen_recent:
            continue
        deduped_parts.append(part)
        if len(seen_recent) >= 24:
            seen_recent.clear()
        seen_recent.add(part)

    has_substantive = any(_is_substantive_ocr_part(part) for part in deduped_parts)
    filtered_parts = [
        part
        for part in deduped_parts
        if _is_substantive_ocr_part(part) or not has_substantive
    ]
    if not filtered_parts:
        filtered_parts = deduped_parts[:]

    short_ratio = sum(1 for part in filtered_parts if len(part) <= 12) / max(len(filtered_parts), 1)
    should_compact = len(filtered_parts) >= 8 and (short_ratio >= 0.55 or quality_score < 0.8)
    if not should_compact:
        return filtered_parts

    compacted: list[str] = []
    buffer: list[str] = []
    for part in filtered_parts:
        candidate = normalize_text(" ".join(buffer + [part]))
        if buffer and (len(candidate) > 36 or len(buffer) >= 3 or len(part) >= 18):
            compacted.append(normalize_text(" ".join(buffer)))
            buffer = [part]
            continue
        buffer.append(part)
        if len(part) >= 18:
            compacted.append(normalize_text(" ".join(buffer)))
            buffer = []
    if buffer:
        compacted.append(normalize_text(" ".join(buffer)))
    return [part for part in compacted if part]


def _is_placeholder_text(text: str) -> bool:
    normalized = normalize_text(text).strip().lower()
    if not normalized:
        return True
    if any(normalized.startswith(prefix) for prefix in PPT_PLACEHOLDER_PREFIXES):
        return True
    return any(pattern.match(normalized) for pattern in PPT_PLACEHOLDER_PATTERNS)


def _sort_slide_elements(elements: list[tuple[int, int, str, str, str, float]]) -> list[tuple[int, int, str, str, str, float]]:
    if len(elements) <= 2:
        return sorted(elements, key=lambda item: (item[0], item[1]))

    sorted_by_left = sorted(elements, key=lambda item: item[1])
    left_positions = [item[1] for item in sorted_by_left]
    min_left = left_positions[0]
    max_left = left_positions[-1]
    horizontal_span = max_left - min_left
    if horizontal_span < 1_600_000:
        return sorted(elements, key=lambda item: (item[0], item[1]))

    column_split = min_left + horizontal_span / 2
    left_column = [item for item in elements if item[1] < column_split]
    right_column = [item for item in elements if item[1] >= column_split]
    if not left_column or not right_column:
        return sorted(elements, key=lambda item: (item[0], item[1]))

    left_span = max(item[1] for item in left_column) - min(item[1] for item in left_column)
    right_span = max(item[1] for item in right_column) - min(item[1] for item in right_column)
    if left_span > 1_600_000 or right_span > 1_600_000:
        return sorted(elements, key=lambda item: (item[0], item[1]))

    left_sorted = sorted(left_column, key=lambda item: (item[0], item[1]))
    right_sorted = sorted(right_column, key=lambda item: (item[0], item[1]))
    left_gap = max(item[0] for item in left_sorted) - min(item[0] for item in left_sorted)
    right_gap = max(item[0] for item in right_sorted) - min(item[0] for item in right_sorted)
    if left_gap < 400_000 or right_gap < 400_000:
        return sorted(elements, key=lambda item: (item[0], item[1]))
    return left_sorted + right_sorted


def _is_ocr_element(element: tuple[int, int, str, str, str, float]) -> bool:
    return element[3] in {"image_ocr", "image_ocr_low_conf"}


def _merge_dense_ocr_elements(elements: list[tuple[int, int, str, str, str, float]]) -> list[tuple[int, int, str, str, str, float]]:
    ocr_elements = [element for element in elements if _is_ocr_element(element)]
    if len(ocr_elements) < 18:
        return elements

    short_ocr_count = sum(1 for _, _, text, _, _, _ in ocr_elements if len(normalize_text(text)) <= 18)
    if short_ocr_count / max(len(ocr_elements), 1) < 0.5:
        return elements

    merged: list[tuple[int, int, str, str, str, float]] = []
    pending_ocr: list[tuple[int, int, str, str, str, float]] = []

    def flush_pending() -> None:
        nonlocal pending_ocr
        if not pending_ocr:
            return
        top = min(item[0] for item in pending_ocr)
        left = min(item[1] for item in pending_ocr)
        quality_score = min(item[5] for item in pending_ocr)
        kind = "image_ocr" if quality_score >= 0.8 else "image_ocr_low_conf"
        content = normalize_text(" / ".join(item[2] for item in pending_ocr if normalize_text(item[2])))
        if content:
            merged.append((top, left, content, kind, "picture-ocr-merged", quality_score))
        pending_ocr = []

    previous_ocr: tuple[int, int, str, str, str, float] | None = None
    for element in elements:
        if not _is_ocr_element(element):
            flush_pending()
            merged.append(element)
            previous_ocr = None
            continue

        top, left, text, _, _, _ = element
        normalized = normalize_text(text)
        if not normalized:
            continue
        if len(normalized) <= 4 and re.fullmatch(r"[A-Za-z0-9]+", normalized):
            continue

        if previous_ocr is None:
            pending_ocr.append(element)
            previous_ocr = element
            continue

        prev_top, prev_left, _, _, _, _ = previous_ocr
        same_row = abs(top - prev_top) <= 260_000
        nearby = abs(left - prev_left) <= 2_400_000
        if same_row and nearby and len(pending_ocr) < 5:
            pending_ocr.append(element)
        else:
            flush_pending()
            pending_ocr.append(element)
        previous_ocr = element

    flush_pending()
    return merged


def _collapse_ocr_heavy_slide(elements: list[tuple[int, int, str, str, str, float]]) -> list[tuple[int, int, str, str, str, float]]:
    ocr_elements = [element for element in elements if _is_ocr_element(element)]
    if len(ocr_elements) < 50:
        return elements

    non_ocr_elements = [element for element in elements if not _is_ocr_element(element)]
    sorted_ocr = sorted(ocr_elements, key=lambda item: (item[0], item[1]))
    merged_ocr: list[tuple[int, int, str, str, str, float]] = []
    buffer: list[tuple[int, int, str, str, str, float]] = []

    def flush_buffer() -> None:
        nonlocal buffer
        if not buffer:
            return
        top = min(item[0] for item in buffer)
        left = min(item[1] for item in buffer)
        quality_score = min(item[5] for item in buffer)
        kind = "image_ocr" if quality_score >= 0.8 else "image_ocr_low_conf"
        content = normalize_text(" / ".join(item[2] for item in buffer if normalize_text(item[2])))
        if content:
            merged_ocr.append((top, left, content, kind, "picture-ocr-collapsed", quality_score))
        buffer = []

    for element in sorted_ocr:
        top, left, text, _, _, _ = element
        normalized = normalize_text(text)
        if not normalized:
            continue
        if not buffer:
            buffer.append(element)
            continue

        prev_top, prev_left, _, _, _, _ = buffer[-1]
        same_band = abs(top - prev_top) <= 360_000
        nearby = abs(left - prev_left) <= 3_400_000
        candidate = normalize_text(" / ".join([item[2] for item in buffer] + [normalized]))
        if (
            len(buffer) >= 8
            or len(candidate) > 180
            or (len(buffer) >= 4 and not same_band)
            or (len(buffer) >= 4 and not nearby)
        ):
            flush_buffer()
        buffer.append(element)
    flush_buffer()

    return non_ocr_elements + merged_ocr


class PptxParser:
    parser_name = "python-pptx"

    def __init__(self, enable_ocr_fallback: bool = True, ocr_language: str = "ch") -> None:
        self.enable_ocr_fallback = enable_ocr_fallback
        self.ocr = PaddleOcrAdapter(language=ocr_language) if enable_ocr_fallback else None

    def parse(self, path: Path) -> ParsedDocument:
        presentation = Presentation(str(path))
        blocks: list[SourceBlock] = []
        markdown_lines: list[str] = []
        warnings: list[str] = []
        ocr_used = False
        title = path.stem

        for index, slide in enumerate(presentation.slides, start=1):
            elements: list[tuple[int, int, str, str, str, float]] = []
            slide_title = ""
            if slide.shapes.title and normalize_text(slide.shapes.title.text):
                slide_title = normalize_text(slide.shapes.title.text)
                if title == path.stem and slide_title:
                    title = slide_title

            for shape in _iter_flattened_shapes(slide.shapes):
                top = _safe_coord(shape, "top")
                left = _safe_coord(shape, "left")
                role = _shape_role(shape, slide_title)
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    lines = []
                    for paragraph in shape.text_frame.paragraphs:
                        text = normalize_text("".join(run.text for run in paragraph.runs) or paragraph.text)
                        if text and not _is_placeholder_text(text):
                            lines.append(text)
                    text = "\n".join(lines)
                    if text:
                        split_parts = _split_shape_text(text)
                        if split_parts:
                            for part_index, part in enumerate(split_parts, start=1):
                                kind = "list" if LIST_MARKER_RE.match(part) else ("title" if role == "title" else "paragraph")
                                label = f"{role}-{part_index}" if len(split_parts) > 1 else role
                                elements.append((top, left, part, kind, label, 1.0))
                if hasattr(shape, "has_table") and shape.has_table:
                    table_md = _table_to_markdown(shape.table)
                    if table_md:
                        elements.append((top, left, table_md, "table", "table", 1.0))
                if (
                    self.enable_ocr_fallback
                    and self.ocr
                    and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                ):
                    try:
                        raw_ocr_text = self.ocr.extract_image_text(shape.image.blob)
                        ocr_text, quality_score = clean_ocr_text(raw_ocr_text)
                        if ocr_text:
                            kind = "image_ocr" if quality_score >= 0.8 else "image_ocr_low_conf"
                            split_parts = _compact_ocr_parts(_split_shape_text(ocr_text) or [ocr_text], quality_score)
                            for part_index, part in enumerate(split_parts, start=1):
                                label = f"picture-ocr-{part_index}" if len(split_parts) > 1 else "picture-ocr"
                                elements.append((top, left, part, kind, label, quality_score))
                            ocr_used = True
                    except Exception as exc:  # pragma: no cover - optional dependency
                        warnings.append(f"第 {index} 页图片 OCR 失败: {exc}")

            if not elements:
                continue

            elements = _merge_dense_ocr_elements(elements)
            elements = _collapse_ocr_heavy_slide(elements)
            elements = _sort_slide_elements(elements)
            slide_heading = slide_title or f"幻灯片 {index}"
            markdown_lines.append(f"## {slide_heading}")
            page_or_slide = f"slide-{index}"
            section_path = f"幻灯片 {index} / {slide_heading}"
            for element_index, element in enumerate(elements, start=1):
                _, _, text, kind, label, quality_score = element
                blocks.append(
                    SourceBlock(
                        page_or_slide=page_or_slide,
                        section_path=f"{section_path} / {label or f'区块 {element_index}'}",
                        content=text,
                        kind=kind,
                        quality_score=quality_score,
                    )
                )
                markdown_lines.append(text)

        return ParsedDocument(
            title=title,
            blocks=blocks,
            raw_markdown="\n\n".join(markdown_lines),
            parser_name=self.parser_name,
            ocr_used=ocr_used,
            warnings=warnings,
        )
